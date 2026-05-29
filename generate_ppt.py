"""
시나브로 팀 - 세법개정 AI 어시스턴트 발표용 PPT 생성 스크립트

실행 방법:
  pip install python-pptx
  python generate_ppt.py

출력: 시나브로_세법개정_AI어시스턴트.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────────────────────────────────────
# 색상 팔레트
# ──────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x64, 0x74, 0x8B)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


def set_slide_bg(slide, color):
    """슬라이드 배경색 설정."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_text(slide, left, top, width, height, text,
                        font_size=14, bold=False, color=DARK_GRAY,
                        fill_color=None, align=PP_ALIGN.LEFT,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """텍스트가 들어간 도형 추가."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_top = Cm(0.3)
    tf.margin_bottom = Cm(0.3)
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.4)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return shape


def add_title_slide(prs):
    """표지 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # 상단 장식 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.3), Inches(10), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # 팀명
    add_shape_with_text(slide, Inches(1), Inches(1.2), Inches(8), Inches(0.7),
                        "팀 시나브로", font_size=18, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # 메인 타이틀
    add_shape_with_text(slide, Inches(0.5), Inches(2.0), Inches(9), Inches(1.5),
                        "세법개정 AI 어시스턴트", font_size=40, bold=True,
                        color=WHITE, align=PP_ALIGN.CENTER)

    # 서브 타이틀
    add_shape_with_text(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.8),
                        "AI Agent 기반 법령 개정 자동화 시스템",
                        font_size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # 구분선
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.5), Inches(3), Pt(2))
    line2.fill.solid()
    line2.fill.fore_color.rgb = ACCENT_BLUE
    line2.line.fill.background()

    # 팀원 정보
    add_shape_with_text(slide, Inches(1.5), Inches(5.0), Inches(7), Inches(1.2),
                        "조장  김의택\n조원  이금석 · 이주호 · 전명수 · 조외영",
                        font_size=14, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # 날짜
    add_shape_with_text(slide, Inches(3), Inches(6.5), Inches(4), Inches(0.5),
                        "2025년 기획재정부 AI 업무혁신 프로젝트",
                        font_size=12, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


def add_toc_slide(prs):
    """목차 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # 왼쪽 색상 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_shape_with_text(slide, Inches(0.7), Inches(0.5), Inches(4), Inches(0.7),
                        "CONTENTS", font_size=28, bold=True, color=NAVY)

    items = [
        ("01", "업무 현황 및 페인포인트"),
        ("02", "AI 서비스 적용 아이디어"),
        ("03", "시스템 구현 개요"),
        ("04", "핵심 기능 시연 흐름"),
        ("05", "기술 아키텍처"),
        ("06", "활용 방안 및 확장 가능성"),
        ("07", "기대 효과"),
    ]

    for i, (num, title) in enumerate(items):
        y = Inches(1.5 + i * 0.8)
        # 번호 원형
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_shape_with_text(slide, Inches(1.8), y, Inches(6), Inches(0.5),
                            title, font_size=16, color=DARK_GRAY)


def add_pain_point_slide(prs):
    """페인포인트 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # 헤더
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = "01  업무 현황 및 페인포인트"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 현황 박스
    add_shape_with_text(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.9),
                        "현황: 법제처 법령안편집기 3.1 활용 중\n"
                        "→ 개정지시문·신구조문대비표 자동 생성만 가능, 나머지는 100% 수작업",
                        font_size=13, color=DARK_GRAY, fill_color=LIGHT_BLUE)

    # 문제점 박스들
    problems = [
        ("❶ 유사 입법례 검색의 어려움", "방대한 법령 DB에서 유사 표현을 직접 찾아 참조해야 함"),
        ("❷ 인용·준용 규정 연쇄 개정 누락 ★", "\"제3항에서 제7항까지\" 등 범위 인용 → Ctrl+F로 탐지 불가"),
        ("❸ 병행 법령 개정 누락 ★", "소득세법↔법인세법 동일 취지 개편 시 소관부서 구분으로 누락 빈번"),
        ("❹ 부칙 효력 판단 오류", "장래효·부진정소급효 구분 미숙 → 부칙 누락 또는 오기 발생"),
    ]

    for i, (title, desc) in enumerate(problems):
        y = Inches(2.5 + i * 1.15)
        # 왼쪽 색상 표시
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Pt(6), Inches(0.9))
        marker.fill.solid()
        marker.fill.fore_color.rgb = RED if "★" in title else ORANGE
        marker.line.fill.background()

        add_shape_with_text(slide, Inches(0.8), y, Inches(8.5), Inches(0.9),
                            f"{title}\n{desc}",
                            font_size=12, color=DARK_GRAY, fill_color=LIGHT_GRAY)

    # 하단 결론
    add_shape_with_text(slide, Inches(0.5), Inches(6.8), Inches(9), Inches(0.5),
                        "⚠ 제도 총괄부처 이미지 악화 + 불필요한 법령 개정 절차 반복 → AI Agent로 해결 가능",
                        font_size=12, bold=True, color=RED)


def add_idea_slide(prs):
    """AI 서비스 적용 아이디어 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = "02  AI 서비스 적용 아이디어"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 입력
    add_shape_with_text(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.6),
                        "💡 입력: 세법개정 요강(개정 핵심 내용)만 텍스트로 입력",
                        font_size=14, bold=True, color=NAVY, fill_color=LIGHT_BLUE)

    # 3단계 흐름
    stages = [
        ("1단계", "개정 조문 초안 자동 생성",
         "• AI가 유사 입법례를 참조하여 조문 구조 자동 작성\n"
         "• 부칙 유형(장래효/부진정소급효 등) 선택 시 AI 자동 출력\n"
         "• 병행 법령(소득세법↔법인세법) 동시 개정안 자동 탐색"),
        ("2단계", "인용·준용 규정 연쇄 확인",
         "• 동일 법령 내, 법·령·칙 간, 타 법령 간 인용 분석\n"
         "• 조문 번호 밀림 영향도 자동 검사\n"
         "• 사용자가 일괄반영 또는 일부반영 선택"),
        ("3단계", "공식 문서 자동 출력",
         "• 개정지시문 + 신구조문대비표 자동 작성\n"
         "• 한글(HWPX) 파일로 즉시 출력\n"
         "• 법령안편집기 3.1과 동일한 포맷 유지"),
    ]

    colors = [ACCENT_BLUE, GREEN, PURPLE]
    for i, (stage, title, desc) in enumerate(stages):
        x = Inches(0.3 + i * 3.2)
        y = Inches(2.2)

        # 단계 뱃지
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.0), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = colors[i]
        badge.line.fill.background()
        btf = badge.text_frame
        btf.paragraphs[0].text = stage
        btf.paragraphs[0].font.size = Pt(11)
        btf.paragraphs[0].font.bold = True
        btf.paragraphs[0].font.color.rgb = WHITE
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 제목
        add_shape_with_text(slide, x, Inches(2.7), Inches(3.0), Inches(0.4),
                            title, font_size=13, bold=True, color=colors[i])

        # 설명 박스
        add_shape_with_text(slide, x, Inches(3.2), Inches(3.0), Inches(2.5),
                            desc, font_size=10, color=DARK_GRAY, fill_color=LIGHT_GRAY)

    # 화살표 연결 표시
    add_shape_with_text(slide, Inches(3.1), Inches(2.3), Inches(0.3), Inches(0.3),
                        "→", font_size=18, bold=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
    add_shape_with_text(slide, Inches(6.3), Inches(2.3), Inches(0.3), Inches(0.3),
                        "→", font_size=18, bold=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # 하단 출력
    add_shape_with_text(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.8),
                        "📄 최종 출력: 법령별 개정안 HWPX 파일 (현행 빨간색 삭제, 개정안 파란색 신설)\n"
                        "     → 법령안편집기 없이도 즉시 활용 가능한 공식 문서 형태",
                        font_size=12, color=NAVY, fill_color=LIGHT_BLUE)


def add_system_overview_slide(prs):
    """시스템 구현 개요 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = "03  시스템 구현 개요"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 핵심 기술 스택
    techs = [
        ("🤖 생성형 AI", "GPT 기반 대규모 언어모델이\n법률 문맥을 이해하고\n개정 초안을 자동 생성"),
        ("🏛 법제처 Open API", "현행 법령 원문을 실시간 조회\n최신 법령 데이터 기반\n정확한 조문 참조 보장"),
        ("🔍 지능형 인용 분석", "정규식+AI 하이브리드로\n인용·준용 규정 자동 탐지\n조문 번호 밀림 연쇄 분석"),
        ("📝 HWPX 문서 생성", "법제처 공식 포맷 준수\n신구조문대비표 자동 편집\n색상 코딩된 변경 표시"),
    ]

    for i, (title, desc) in enumerate(techs):
        x = Inches(0.3 + i * 2.4)
        y = Inches(1.3)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(2.8))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(1)

        add_shape_with_text(slide, x + Inches(0.1), y + Inches(0.2), Inches(2.0), Inches(0.5),
                            title, font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_shape_with_text(slide, x + Inches(0.1), y + Inches(0.8), Inches(2.0), Inches(1.8),
                            desc, font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # 워크플로우 요약
    add_shape_with_text(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(0.5),
                        "▎ End-to-End 자동화 파이프라인",
                        font_size=14, bold=True, color=NAVY)

    flow_items = [
        ("법령 검색\n& 조문 선택", ACCENT_BLUE),
        ("개정 요강\n입력", ACCENT_BLUE),
        ("AI 초안\n생성", GREEN),
        ("인용규정\n연쇄 분석", ORANGE),
        ("병행법령\n자동 탐지", PURPLE),
        ("HWPX\n출력", RED),
    ]

    for i, (text, color) in enumerate(flow_items):
        x = Inches(0.4 + i * 1.55)
        y = Inches(5.2)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.3), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        btf = box.text_frame
        btf.word_wrap = True
        btf.paragraphs[0].text = text
        btf.paragraphs[0].font.size = Pt(9)
        btf.paragraphs[0].font.bold = True
        btf.paragraphs[0].font.color.rgb = WHITE
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE

        if i < len(flow_items) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           x + Inches(1.35), y + Inches(0.4),
                                           Inches(0.2), Pt(3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MEDIUM_GRAY
            arrow.line.fill.background()

    # 하단 강조
    add_shape_with_text(slide, Inches(0.5), Inches(6.6), Inches(9), Inches(0.5),
                        "✓ 사용자는 '무엇을 개정할지'만 입력 → AI가 '어떻게 개정할지' 전 과정 지원",
                        font_size=12, bold=True, color=ACCENT_BLUE)


def add_demo_flow_slide(prs):
    """핵심 기능 시연 흐름 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = "04  핵심 기능 시연 흐름"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    steps = [
        ("STEP 1", "법령 검색 & 조문 선택",
         "법령명을 입력하면 법제처 Open API로\n실시간 검색. 개정할 조문을 직접 선택.",
         ACCENT_BLUE),
        ("STEP 2", "개정 요강 입력 & AI 초안 생성",
         "\"법인세율 20%→19% 인하\" 등 핵심만 입력.\nAI가 유사 입법례를 참조해 조문 초안 즉시 생성.\n부칙 유형 선택 시 적용례까지 자동 작성.",
         GREEN),
        ("STEP 3", "인용 규정 연쇄 분석",
         "개정 조문을 인용하는 모든 규정 자동 탐지.\n\"제3항~제7항\" 범위 인용도 정밀 분석.\n조문 번호 밀림 영향까지 사전 경고.",
         ORANGE),
        ("STEP 4", "병행 법령 동시 개정",
         "소득세법↔법인세법 등 동일 취지 제도를\nAI가 자동으로 크로스체크.\n누락 없는 일괄 개정안 생성.",
         PURPLE),
        ("STEP 5", "HWPX 공식 문서 출력",
         "신구조문대비표를 법제처 포맷 그대로\nHWPX 파일로 즉시 생성.\n현행(빨강)·개정안(파랑) 색상 구분 표시.",
         RED),
    ]

    for i, (step, title, desc, color) in enumerate(steps):
        y = Inches(1.2 + i * 1.2)

        # 스텝 뱃지
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.4), y, Inches(0.9), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        btf = badge.text_frame
        btf.paragraphs[0].text = step
        btf.paragraphs[0].font.size = Pt(9)
        btf.paragraphs[0].font.bold = True
        btf.paragraphs[0].font.color.rgb = WHITE
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 제목
        add_shape_with_text(slide, Inches(1.5), y, Inches(3.0), Inches(0.35),
                            title, font_size=12, bold=True, color=color)

        # 설명
        add_shape_with_text(slide, Inches(4.5), y - Inches(0.1), Inches(5.2), Inches(1.0),
                            desc, font_size=10, color=DARK_GRAY, fill_color=LIGHT_GRAY)


def add_architecture_slide(prs):
    """기술 아키텍처 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = "05  기술 아키텍처"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 계층 구조
    layers = [
        ("사용자 인터페이스 계층", "웹 기반 인터랙티브 UI · 실시간 미리보기 · 단계별 가이드",
         ACCENT_BLUE, Inches(1.2)),
        ("AI 에이전트 계층", "대규모 언어모델(LLM) · 법률 도메인 특화 프롬프트 · 유사 입법례 추론",
         GREEN, Inches(2.6)),
        ("법령 분석 엔진", "인용 규정 정규식 파서 · 조문 번호 밀림 분석 · 병행 법령 의미 매칭",
         ORANGE, Inches(4.0)),
        ("데이터 소스 계층", "법제처 Open API (실시간) · 세법 동의어 사전 · 부칙 템플릿 라이브러리",
         PURPLE, Inches(5.4)),
    ]

    for title, desc, color, y in layers:
        # 메인 박스
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.8), y, Inches(8.4), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        add_shape_with_text(slide, Inches(1.0), y + Inches(0.1), Inches(3.5), Inches(0.4),
                            title, font_size=13, bold=True, color=WHITE)
        add_shape_with_text(slide, Inches(1.0), y + Inches(0.5), Inches(7.8), Inches(0.5),
                            desc, font_size=11, color=WHITE)

    # 오른쪽 세로 텍스트 - 출력 포맷
    add_shape_with_text(slide, Inches(0.8), Inches(6.7), Inches(8.4), Inches(0.5),
                        "📎 출력 포맷: HWPX (한글 공식) · 신구조문대비표 · 개정지시문 · 부칙",
                        font_size=11, color=NAVY, fill_color=LIGHT_BLUE)


def add_usage_expansion_slide(prs):
    """활용 방안 및 확장 가능성 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = "06  활용 방안 및 확장 가능성"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 현재 활용 방안
    add_shape_with_text(slide, Inches(0.5), Inches(1.2), Inches(4.3), Inches(0.5),
                        "▎ 현재 활용 방안 (세제실)", font_size=14, bold=True, color=NAVY)

    current_uses = [
        "• 정기 세법개정 시 개정안 초안 일괄 생성",
        "• 인용 규정 누락 사전 방지 (Zero-Miss 목표)",
        "• 소득세법↔법인세법 병행 개정 자동 검증",
        "• 신규 담당자 부칙 작성 가이드 역할",
        "• 법령안편집기 로딩 오류 회피 (독립 시스템)",
    ]

    for i, text in enumerate(current_uses):
        add_shape_with_text(slide, Inches(0.7), Inches(1.8 + i * 0.5), Inches(4.0), Inches(0.4),
                            text, font_size=11, color=DARK_GRAY)

    # 확장 가능성
    add_shape_with_text(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.5),
                        "▎ 확장 가능성", font_size=14, bold=True, color=PURPLE)

    expansions = [
        ("🏦 금융위", "자본시장법, 보험업법 등"),
        ("🏗 국토부", "건축법, 주택법 시행령"),
        ("👨‍⚕ 복지부", "국민건강보험법, 의료법"),
        ("⚖ 법무부", "형법, 민법 개정"),
        ("🌐 범정부", "모든 법률·시행령·시행규칙"),
    ]

    for i, (icon, desc) in enumerate(expansions):
        y = Inches(1.8 + i * 0.7)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(5.3), y, Inches(4.0), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()

        add_shape_with_text(slide, Inches(5.4), y + Inches(0.05), Inches(3.8), Inches(0.45),
                            f"{icon}  {desc}", font_size=11, color=DARK_GRAY)

    # 하단 로드맵
    add_shape_with_text(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.4),
                        "▎ 미래 확장 로드맵", font_size=14, bold=True, color=NAVY)

    roadmap = [
        ("Phase 1\n(현재)", "세법 특화\n프로토타입", ACCENT_BLUE),
        ("Phase 2", "법·령·칙\n전체 연동", GREEN),
        ("Phase 3", "타 부처\n법령 확장", ORANGE),
        ("Phase 4", "범정부\n입법 플랫폼", PURPLE),
    ]

    for i, (phase, desc, color) in enumerate(roadmap):
        x = Inches(0.6 + i * 2.3)
        y = Inches(6.1)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.0), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        add_shape_with_text(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.8), Inches(0.4),
                            phase, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_shape_with_text(slide, x + Inches(0.1), y + Inches(0.45), Inches(1.8), Inches(0.45),
                            desc, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

        if i < len(roadmap) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           x + Inches(2.05), y + Inches(0.45),
                                           Inches(0.2), Pt(3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MEDIUM_GRAY
            arrow.line.fill.background()


def add_expected_effects_slide(prs):
    """기대 효과 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = "07  기대 효과"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 정량적 효과
    quant_effects = [
        ("⏱", "업무 시간", "70% 단축", "법령 개정안 초안 작성\n수일 → 수 분"),
        ("🎯", "정확도", "인용 누락 0건", "AI 전수 검사로\n연쇄 개정 누락 방지"),
        ("⚡", "처리 속도", "실시간", "법령 로딩 오류 제거\n즉시 조문 조회"),
        ("📊", "커버리지", "6개 법령 쌍", "병행 법령 자동 탐지\n확장 구조 확보"),
    ]

    for i, (icon, label, value, desc) in enumerate(quant_effects):
        x = Inches(0.3 + i * 2.4)
        y = Inches(1.3)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(1)

        add_shape_with_text(slide, x + Inches(0.1), y + Inches(0.15), Inches(2.0), Inches(0.4),
                            f"{icon} {label}", font_size=11, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
        add_shape_with_text(slide, x + Inches(0.1), y + Inches(0.6), Inches(2.0), Inches(0.5),
                            value, font_size=16, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
        add_shape_with_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.0), Inches(0.9),
                            desc, font_size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # 정성적 효과
    add_shape_with_text(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(0.4),
                        "▎ 정성적 기대 효과", font_size=14, bold=True, color=NAVY)

    qual_effects = [
        "✅ 제도 총괄부처로서의 신뢰도 회복 — 개정 누락·오기로 인한 이미지 악화 원천 차단",
        "✅ 불필요한 법령 개정 절차 반복 제거 — 보완 개정 건수 대폭 감소",
        "✅ 신규 담당자 온보딩 시간 단축 — AI가 부칙 유형·인용 규정 가이드 역할 수행",
        "✅ 정기 세법개정 시즌 업무 부하 획기적 경감 — 핵심 정책 검토에 집중 가능",
        "✅ 디지털 전환 선도 부처 위상 확립 — AI 기반 입법 지원의 정부 최초 사례",
    ]

    for i, text in enumerate(qual_effects):
        add_shape_with_text(slide, Inches(0.7), Inches(4.5 + i * 0.5), Inches(8.5), Inches(0.4),
                            text, font_size=11, color=DARK_GRAY)


def add_closing_slide(prs):
    """마무리 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # 상단 장식
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.3), Inches(10), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_shape_with_text(slide, Inches(1), Inches(2.0), Inches(8), Inches(1.0),
                        "\"AI가 법령을 이해하고,\n사람은 정책에 집중하는 시대\"",
                        font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_shape_with_text(slide, Inches(1.5), Inches(3.5), Inches(7), Inches(0.8),
                        "세법개정의 패러다임을 바꾸는 첫 걸음",
                        font_size=18, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # 구분선
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.5), Inches(3), Pt(2))
    line2.fill.solid()
    line2.fill.fore_color.rgb = ACCENT_BLUE
    line2.line.fill.background()

    add_shape_with_text(slide, Inches(2), Inches(5.0), Inches(6), Inches(1.0),
                        "팀 시나브로\n김의택 · 이금석 · 이주호 · 전명수 · 조외영",
                        font_size=14, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    add_shape_with_text(slide, Inches(2.5), Inches(6.3), Inches(5), Inches(0.5),
                        "감사합니다", font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def main():
    """PPT 생성 메인 함수."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_toc_slide(prs)
    add_pain_point_slide(prs)
    add_idea_slide(prs)
    add_system_overview_slide(prs)
    add_demo_flow_slide(prs)
    add_architecture_slide(prs)
    add_usage_expansion_slide(prs)
    add_expected_effects_slide(prs)
    add_closing_slide(prs)

    output_path = os.path.join(os.path.dirname(__file__), "시나브로_세법개정_AI어시스턴트.pptx")
    prs.save(output_path)
    print(f"✅ PPT 생성 완료: {output_path}")
    print(f"   총 {len(prs.slides)}장 슬라이드")


if __name__ == "__main__":
    main()
